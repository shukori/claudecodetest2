# -*- coding: utf-8 -*-
"""HYUGA PRIMARY CARE 在宅訪問薬局 競合マップ 3枚（福岡・九州・関東）"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ===== カラーパレット =====
NAVY   = RGBColor(0x1F, 0x3A, 0x5F)   # タイトル帯
BLUE   = RGBColor(0x2E, 0x6D, 0xA4)   # 見出し
LBLUE  = RGBColor(0xE8, 0xF1, 0xF8)   # 行うす青
GREEN  = RGBColor(0x2E, 0x8B, 0x57)   # 自社強調
LGREEN = RGBColor(0xE3, 0xF2, 0xE9)
GRAY   = RGBColor(0x5A, 0x5A, 0x5A)
LGRAY  = RGBColor(0xF2, 0xF2, 0xF2)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x22, 0x22, 0x22)
RED    = RGBColor(0xC0, 0x39, 0x2B)
FONT   = "IPAGothic"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
BLANK = prs.slide_layouts[6]


def set_run(run, size, color, bold=False):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = FONT


def add_text(slide, l, t, w, h, text, size, color, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, fill=None,
             line=None, wrap=True):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(5); tf.margin_right = Pt(5)
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        set_run(r, size, color, bold)
    if fill is not None:
        box.fill.solid(); box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    if line is not None:
        box.line.color.rgb = line; box.line.width = Pt(0.75)
    else:
        box.line.fill.background()
    return box


def cell(slide, l, t, w, h, text, size, color, fill, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = WHITE; box.line.width = Pt(1.2)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(6); tf.margin_right = Pt(6)
    tf.margin_top = Pt(3); tf.margin_bottom = Pt(3)
    for i, ln in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = 1.0
        r = p.add_run(); r.text = ln
        set_run(r, size, color, bold)
    return box


def header(slide, region, position):
    # タイトル帯
    bar = slide.shapes.add_textbox(0, 0, SW, Inches(0.95))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    tf = bar.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(18)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = f"HYUGA PRIMARY CARE（きらり薬局）── {region} 競合マップ"
    set_run(r, 22, WHITE, True)
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = f"ポジション：{position}"
    set_run(r2, 12.5, RGBColor(0xCF, 0xE2, 0xF3), False)


def table(slide, top, headers, col_w, rows, row_h_list):
    """rows: list of (cells:list, fill, textcolor, bold). 各セル内改行可"""
    left0 = Inches(0.35)
    # header row
    x = left0
    hh = Inches(0.42)
    for j, htxt in enumerate(headers):
        cell(slide, x, top, col_w[j], hh, htxt, 12, WHITE, BLUE, bold=True,
             align=PP_ALIGN.CENTER)
        x += col_w[j]
    y = top + hh
    for ridx, (cells, fill, tcolor, bold) in enumerate(rows):
        rh = row_h_list[ridx]
        x = left0
        for j, ctxt in enumerate(cells):
            al = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.LEFT
            cell(slide, x, y, col_w[j], rh, ctxt,
                 10.5 if j > 0 else 11, tcolor, fill, bold=bold, align=al)
            x += col_w[j]
        y += rh
    return y


def summary(slide, top, lines):
    box = slide.shapes.add_textbox(Inches(0.35), top, Inches(12.63), Inches(0.95))
    box.fill.solid(); box.fill.fore_color.rgb = LGREEN
    box.line.color.rgb = GREEN; box.line.width = Pt(1.5)
    tf = box.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(12); tf.margin_top = Pt(4); tf.margin_bottom = Pt(4)
    for i, (label, txt, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.05
        if label:
            r = p.add_run(); r.text = label
            set_run(r, 11.5, GREEN, True)
        r2 = p.add_run(); r2.text = txt
        set_run(r2, 11.5, DARK, bold)


# =========================================================
# スライド①：福岡
# =========================================================
s = prs.slides.add_slide(BLANK)
header(s, "福岡県", "「在宅に振り切った唯一の上場プレイヤー」── 福岡が根拠地（先行者）")
colw = [Inches(2.3), Inches(2.5), Inches(3.0), Inches(4.83)]
heads = ["競合", "タイプ・規模", "福岡での在宅姿勢", "HYUGAとの戦い方"]
rows = [
    (["そうごう薬局\n(総合メディカル)", "大手チェーン\n福岡72店超／県内最大級", "在宅も対応するが\n主軸は門前・外来", "同一施設で両者が営業\n→施設リレーションの深さで棲み分け"], WHITE, DARK, False),
    (["アイン薬局\n(アインHD)", "全国大手\n全国836店超／在宅90%実施", "福岡県内に複数店\n在宅実績あるが外来主体", "規模では劣後。在宅『本気度』と\n24時間対応で差別化"], LGRAY, DARK, False),
    (["日本調剤", "全国大手\n全国697店(2022)", "福岡に複数店\n門前・外来中心", "在宅比率が低く施設在宅では\n直接競合になりにくい"], WHITE, DARK, False),
    (["さくら薬局/新生堂\nタケシタ調剤 等", "地場チェーン\n北部九州中心", "部分的に在宅\n外来が中心", "施設・ケアマネ介入を先着で固め\nKIRARI PRIMEで連携先化"], LGRAY, DARK, False),
    (["★ きらり薬局\n(HYUGA)", "在宅特化・上場\n福岡中心37店／施設隣接立地", "在宅売上60%／在宅9,000名\n1店=門前の約7倍患者", "施設隣接立地＋KIRARI PRIMEで\n地場薬局を競合から味方へ転換"], LGREEN, GREEN, True),
]
rhs = [Inches(0.78), Inches(0.78), Inches(0.70), Inches(0.78), Inches(0.90)]
yend = table(s, Inches(1.15), heads, colw, rows, rhs)
summary(s, yend + Inches(0.10), [
    ("戦い方の要約：", "大手とは「在宅への本気度（売上比率・24h・施設ノウハウ）」で差別化。", True),
    ("", "地場とは「施設囲い込みの先着」と「KIRARI PRIME（競合→連携化）」で無力化する。", False),
])

# =========================================================
# スライド②：九州
# =========================================================
s = prs.slides.add_slide(BLANK)
header(s, "九州全域", "「福岡ドミナント確立済み・他県は未開拓の空白地帯」")
colw = [Inches(2.3), Inches(2.5), Inches(3.0), Inches(4.83)]
heads = ["競合", "タイプ・規模", "九州での動向", "脅威度と対策"]
rows = [
    (["アインHD", "全国大手\n全国1,000店超／在宅90%", "2021年 エス・ケー・\nファーマシーと業務資本提携", "★最大脅威。九州に足場を作り水面下で\n在宅強化中。先にドミナント堀を深める"], WHITE, DARK, False),
    (["そうごう薬局\n(総合メディカル)", "大手チェーン\n全国742店", "福岡を拠点に九州各県へ\n積極出店", "主戦場は外来。在宅は付帯水準だが\nM&Aで在宅強化の可能性あり"], LGRAY, DARK, False),
    (["コスモス薬品", "ドラッグストア\n九州発祥・大量出店", "九州全土に展開\n(調剤併設型)", "現時点は外来・OTC主体。\n在宅の直接競合には今のところならず"], WHITE, DARK, False),
    (["各県地場チェーン", "熊本・鹿児島・宮崎\n長崎等の県内中小", "県内シェア保有\n在宅は局所的", "在宅は小規模。施設紹介ネットを持つ\n場合は KIRARI PRIME で取込"], LGRAY, DARK, False),
    (["★ きらり薬局\n(HYUGA)", "在宅特化・上場\n西日本37店", "実質「福岡＋佐賀」\nに留まる", "福岡以外の九州は未進出が大半\n＝チャンス兼死角"], LGREEN, GREEN, True),
]
rhs = [Inches(0.82), Inches(0.78), Inches(0.78), Inches(0.78), Inches(0.78)]
yend = table(s, Inches(1.15), heads, colw, rows, rhs)
summary(s, yend + Inches(0.10), [
    ("戦い方の要約：", "「九州制覇」より「福岡密度を高めモデルを完成→次の1〜2県でM&A先行」が現実解。", True),
    ("", "大手（特にアインHD）のM&A攻勢より先に、熊本・鹿児島の在宅特化薬局（売り案件）を押さえる。", False),
])

# =========================================================
# スライド③：関東（千葉・神奈川・東京）
# =========================================================
s = prs.slides.add_slide(BLANK)
header(s, "関東（千葉・神奈川・東京）", "「福岡実績を持ち込んだ挑戦者。施設在宅は差別化可だが、ブランド・採用で地力差」")
colw = [Inches(2.3), Inches(2.5), Inches(3.0), Inches(4.83)]
heads = ["競合", "タイプ・規模", "関東での在宅", "HYUGAとの戦い方"]
rows = [
    (["ヤックス\n(千葉薬品)", "千葉地場大手\n千葉全域に展開", "千葉県内で在宅医療に\n積極対応", "★最大の直接競合。地域関係で先行。\n施設在宅にどこまで本腰か要確認"], WHITE, DARK, False),
    (["アイセイ薬局", "大手チェーン\n全国301店／関東主力", "薬剤師訪問サービスあり\nただし外来中心", "在宅売上比率はHYUGAと比較に\nならない水準。施設密着で差別化"], LGRAY, DARK, False),
    (["日本調剤", "全国大手\n全国697店／関東が主力", "関東主力市場\n在宅も対応するが外来主戦場", "2025年TOB等 業界再編の中心。\n規模で圧倒も在宅特化ではない"], WHITE, DARK, False),
    (["アイン薬局\n(アインHD)", "全国大手\n全国1,000店超", "関東全域に大量出店\n在宅も実施", "在宅実績は広いが『付帯』止まり。\n施設密着モデルではHYUGAに分"], LGRAY, DARK, False),
    (["★ きらり薬局\n(HYUGA)", "在宅特化・上場\n関東17店(拡張中)", "千葉9・神奈川6・東京2店\n2025/6 本千葉店開局", "まだ「面」未完成。施設関係構築が急務\nKIRARI PRIMEで地場中小を取込"], LGREEN, GREEN, True),
]
rhs = [Inches(0.80), Inches(0.78), Inches(0.78), Inches(0.78), Inches(0.82)]
yend = table(s, Inches(1.15), heads, colw, rows, rhs)
summary(s, yend + Inches(0.10), [
    ("関東の難所：", "競合が多く市場成熟。福岡型の『空白地帯×先行者優位』は通用しない。採用コストも高い。", True),
    ("戦い方：", "福岡型『施設隣接×24h在宅』を持込み千葉・神奈川を面で固め、KIRARI PRIMEで資産ライト拡大。", False),
])

out = "/home/user/ClaudeCodetest2/HYUGA_在宅訪問薬局_競合マップ.pptx"
prs.save(out)
print("saved:", out)
